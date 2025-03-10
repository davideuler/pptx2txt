import codecs
import collections.abc
from distutils.errors import LinkError  # noqa: F401
import os.path
import sys
from datetime import datetime

import requests

from pptx import Presentation
from pptx.shapes.group import GroupShape

version = '2021-12-20'

# spell check from: https://github.com/shibing624/pycorrector.git
# pycorrector/examples/flask_server_demo.py
corrector_url = "http://127.0.0.1:5000/macbert_correct"

def log(message):
    timestamp = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
    print(f'{timestamp} {message}', file=sys.stderr)


def treat_shape(shape):
    lines = []
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            stripped = paragraph.text.strip()
            if stripped:
                lines.append(stripped)
    elif shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                stripped = cell.text.strip()
                if stripped:
                    lines.append(stripped)
    elif isinstance(shape, GroupShape):
        for item in shape.shapes:
            lines += treat_shape(item)
    return lines


log(f'ppt2txt - version {version} by Shinichi Akiyama')

for file in sys.argv[1:]:
    base, ext = os.path.splitext(file)
    textfile = f'{base}.txt'

    presentation = Presentation(file)
    log(f'{file} was opened.')

    lines = []
    for i, slide in enumerate(presentation.slides):
        lines.append(f'--- Slide {i + 1} ---')
        for shape in slide.shapes:
            lines += treat_shape(shape)

    with codecs.open(textfile, 'w', 'utf-8') as f:
        for line in lines:
            print(line, file=f)

            ##  call spell checker service, and print the correction suggest
            if '--- Slide ' in line:
                print(line)
            else:
                r = requests.post(corrector_url, json={'text': line})
                if not '[]' in r.text: # when result contains diff result
                    print("===== correct: %s" % r.text)
    log(f'{textfile} was saved.')

log('All files were processed.')
